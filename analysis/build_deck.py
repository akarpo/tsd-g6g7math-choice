#!/usr/bin/env python3
"""
build_deck.py - Generate Troy School District G6-G7 Math Executive Summary deck.

Produces a polished 22-slide PowerPoint at:
  deck/Troy_G6G7_Math_Executive_Summary.pptx

Visual style matched to the Troy SD K-5 ELA reference presentation:
  - Light gray slide backgrounds (#EAEDF2)
  - Dark navy header bar with white title + lighter subtitle
  - Dark navy footer bar with repo link + slide counter
  - Calibri throughout, maximised chart/panel sizes, minimal whitespace

Requires: python-pptx
"""

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -- paths -------------------------------------------------------------------
PROJECT = Path(__file__).resolve().parent.parent
CHARTS  = PROJECT / "charts"
DECK    = PROJECT / "deck"
OUTPUT  = DECK / "Troy_G6G7_Math_Executive_Summary.pptx"

DECK.mkdir(parents=True, exist_ok=True)

# -- colour palette -----------------------------------------------------------
TROY_BLUE      = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT_RED     = RGBColor(0xC8, 0x30, 0x2F)
ACCENT_GREEN   = RGBColor(0x1F, 0x7A, 0x3D)
ACCENT_ORANGE  = RGBColor(0xB7, 0x79, 0x1F)
LIGHT_GREEN    = RGBColor(0xE8, 0xF5, 0xEE)
LIGHT_ORANGE   = RGBColor(0xFF, 0xF4, 0xE0)
LIGHT_RED      = RGBColor(0xFB, 0xE7, 0xE6)
GRAY_BG        = RGBColor(0xEA, 0xED, 0xF2)   # slide background
GRAY_LIGHT     = RGBColor(0xF2, 0xF4, 0xF7)
GRAY_DARK      = RGBColor(0x33, 0x33, 0x33)
GRAY_MID       = RGBColor(0x77, 0x77, 0x77)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
SUBTITLE_CLR   = RGBColor(0xA0, 0xB4, 0xCC)
DARK_GREEN_HDR = RGBColor(0x1F, 0x7A, 0x3D)
DARK_RED_HDR   = RGBColor(0xA8, 0x20, 0x1A)

# -- slide geometry -----------------------------------------------------------
SLIDE_W      = Inches(13.333)
SLIDE_H      = Inches(7.5)
HEADER_H     = Inches(0.95)
FOOTER_H     = Inches(0.28)
FOOTER_Y     = SLIDE_H - FOOTER_H           # 7.22"
CONTENT_TOP  = HEADER_H                      # 0.95" — no gap
CONTENT_BOT  = FOOTER_Y                      # 7.22"
CONTENT_H    = CONTENT_BOT - CONTENT_TOP     # ~6.27"
TOTAL_SLIDES = 22

# Chart + panel standard positions
CHART_L  = Inches(0.2)
CHART_T  = Inches(1.0)
CHART_W  = Inches(8.6)
CHART_H  = Inches(5.9)
PANEL_L  = Inches(9.0)
PANEL_T  = Inches(1.0)
PANEL_W  = Inches(4.1)
PANEL_H  = Inches(5.9)
PANEL_PAD = Inches(0.15)

# -- presentation scaffold ----------------------------------------------------
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK_LAYOUT = prs.slide_layouts[6]  # blank


# =============================================================================
# Helper functions
# =============================================================================

def add_rect(slide, left, top, width, height, color):
    """Add a solid-color filled rectangle shape (no border)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=14,
             bold=False, color=GRAY_DARK, alignment='left', font_name='Calibri',
             valign='top'):
    """Add a single-paragraph text box with word wrap."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if valign == 'middle':
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == 'bottom':
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align_map.get(alignment, PP_ALIGN.LEFT)
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=14,
                  bold=False, color=GRAY_DARK, alignment='left',
                  font_name='Calibri', line_spacing=None, valign='top',
                  space_after=None):
    """Add a text box with multiple paragraphs (one per item in *lines*)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    va = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = va.get(valign, MSO_ANCHOR.TOP)
    align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align_map.get(alignment, PP_ALIGN.LEFT)
        if space_after is not None:
            p.space_after = Pt(space_after)
        else:
            p.space_after = Pt(font_size * 0.3)
    return txBox


def add_rich_text(slide, left, top, width, height, runs_list,
                  alignment='left', valign='top'):
    """Textbox with multiple paragraphs each containing multiple styled runs.

    runs_list: [ [ {text, font_size, bold, color, font_name}, ... ], ... ]
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    va = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = va.get(valign, MSO_ANCHOR.TOP)
    align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}

    for i, para_runs in enumerate(runs_list):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map.get(alignment, PP_ALIGN.LEFT)
        p.space_after = Pt(4)
        for j, rd in enumerate(para_runs):
            if j == 0 and not p.runs:
                run = p.add_run()
            else:
                run = p.add_run()
            run.text = rd.get('text', '')
            run.font.size = Pt(rd.get('font_size', 14))
            run.font.bold = rd.get('bold', False)
            run.font.color.rgb = rd.get('color', GRAY_DARK)
            run.font.name = rd.get('font_name', 'Calibri')
    return txBox


def set_slide_bg(slide, color):
    """Set the slide background to a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def header_bar(slide, title, subtitle=None):
    """Dark navy header bar, 0.95" tall, with white title and lighter subtitle."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, HEADER_H, TROY_BLUE)
    add_text(slide, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.5),
             title, font_size=23, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.52), Inches(12.3), Inches(0.35),
                 subtitle, font_size=12, color=SUBTITLE_CLR)


def footer_bar(slide, slide_num):
    """Dark navy footer bar at bottom with left text + right slide counter."""
    add_rect(slide, Inches(0), FOOTER_Y, SLIDE_W, FOOTER_H, TROY_BLUE)
    add_text(slide, Inches(0.4), FOOTER_Y + Inches(0.04), Inches(9), Inches(0.22),
             "Troy SD G6-G7 Math — Executive Summary • github.com/akarpo/tsd-g6g7math-choice",
             font_size=8, color=WHITE)
    add_text(slide, Inches(11.5), FOOTER_Y + Inches(0.04), Inches(1.5), Inches(0.22),
             f"{slide_num} / {TOTAL_SLIDES}",
             font_size=8, color=WHITE, alignment='right')


def new_slide(title, subtitle=None, slide_num=1, bg_color=GRAY_BG):
    """Create a new blank slide with standard header, footer, and background."""
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_slide_bg(slide, bg_color)
    header_bar(slide, title, subtitle)
    footer_bar(slide, slide_num)
    return slide


def embed_chart(slide, chart_name, left=None, top=None, width=None, height=None):
    """Embed a chart PNG if it exists, otherwise show a placeholder."""
    left = left or CHART_L
    top = top or CHART_T
    width = width or CHART_W
    height = height or CHART_H
    chart_path = CHARTS / chart_name
    if chart_path.exists():
        slide.shapes.add_picture(str(chart_path), left, top, width, height)
    else:
        add_rect(slide, left, top, width, height, GRAY_LIGHT)
        add_text(slide, left, top, width, height,
                 f"[Chart: {chart_name}]",
                 font_size=12, color=GRAY_MID, alignment='center', valign='middle')


def accent_bar(slide, left, top, width, color):
    """Thin 3px horizontal accent bar."""
    add_rect(slide, left, top, width, Inches(0.04), color)


def panel_box(slide, left, top, width, height, bg_color):
    """Draw a panel background and return coordinates for inner content."""
    add_rect(slide, left, top, width, height, bg_color)
    return (left + PANEL_PAD, top + PANEL_PAD,
            width - 2 * PANEL_PAD, height - 2 * PANEL_PAD)


# =============================================================================
# Slide 1 — Title
# =============================================================================
def build_slide_01():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_slide_bg(slide, TROY_BLUE)
    # Full TROY_BLUE background rect to ensure complete coverage
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, TROY_BLUE)

    # Main title — left-aligned
    add_text(slide, Inches(1.0), Inches(1.6), Inches(11), Inches(0.8),
             "Troy SD G6-G7 Math", font_size=36, bold=True, color=WHITE)

    # Subtitle
    add_text(slide, Inches(1.0), Inches(2.5), Inches(11), Inches(0.7),
             "Has IM + Detracking Outperformed Its Characteristics?",
             font_size=20, color=WHITE)

    # Thin white horizontal line
    add_rect(slide, Inches(1.0), Inches(3.35), Inches(6.0), Inches(0.02), WHITE)

    # Sub-subtitle
    add_text(slide, Inches(1.0), Inches(3.6), Inches(11), Inches(0.5),
             "Executive Summary — SEDA 2025.1 + M-STEP Data",
             font_size=14, color=WHITE)

    # Bottom info lines
    add_text(slide, Inches(1.0), Inches(5.8), Inches(11), Inches(0.4),
             "Prepared May 2026 • 296-district national benchmark • 2009–2025 coverage",
             font_size=11, color=SUBTITLE_CLR)
    add_text(slide, Inches(1.0), Inches(6.2), Inches(11), Inches(0.4),
             "Repo: github.com/akarpo/tsd-g6g7math-choice",
             font_size=11, color=SUBTITLE_CLR)

    # Footer bar
    footer_bar(slide, 1)


# =============================================================================
# Slide 2 — The Question and the One-Line Answer
# =============================================================================
def build_slide_02():
    slide = new_slide(
        "The Question — and the One-Line Answer",
        "Three independent lines of evidence converge on the same answer",
        slide_num=2)

    # QUESTION — left block
    qx, qy, qw, qh = Inches(0.4), Inches(1.15), Inches(4.0), Inches(2.6)
    add_rect(slide, qx, qy, qw, qh, WHITE)
    add_text(slide, qx + Inches(0.2), qy + Inches(0.12), Inches(2), Inches(0.35),
             "QUESTION", font_size=13, bold=True, color=TROY_BLUE)
    add_text(slide, qx + Inches(0.2), qy + Inches(0.55), qw - Inches(0.4), Inches(1.8),
             "Has Troy’s Fall 2023 adoption of Illustrative Mathematics "
             "and detracking of G6+G7 math produced results that outperform "
             "its characteristics?",
             font_size=15, color=GRAY_DARK)

    # ANSWER — center block
    ax, ay, aw, ah = Inches(4.6), Inches(1.15), Inches(4.2), Inches(2.6)
    add_rect(slide, ax, ay, aw, ah, LIGHT_RED)
    add_text(slide, ax + Inches(0.2), ay + Inches(0.12), Inches(2), Inches(0.35),
             "ANSWER", font_size=13, bold=True, color=ACCENT_RED)
    add_text(slide, ax + Inches(0.2), ay + Inches(0.55), Inches(1.5), Inches(0.6),
             "No.", font_size=32, bold=True, color=ACCENT_RED)
    add_text(slide, ax + Inches(0.2), ay + Inches(1.25), aw - Inches(0.4), Inches(1.2),
             "Troy ranks 197 of 295 level-matched peers (bottom third) "
             "in math recovery during the exact IM adoption window.",
             font_size=14, color=GRAY_DARK)

    # Three evidence boxes at bottom
    box_y = Inches(4.1)
    box_h = Inches(2.9)
    box_w = Inches(4.0)
    gap = Inches(0.15)

    # Box 1 — LIGHT_RED
    b1x = Inches(0.4)
    add_rect(slide, b1x, box_y, box_w, box_h, LIGHT_RED)
    accent_bar(slide, b1x, box_y, box_w, ACCENT_RED)
    add_text(slide, b1x + Inches(0.15), box_y + Inches(0.15), box_w - Inches(0.3), Inches(0.6),
             "197 / 295", font_size=30, bold=True, color=ACCENT_RED)
    add_text(slide, b1x + Inches(0.15), box_y + Inches(0.85), box_w - Inches(0.3), Inches(1.8),
             "Level-matched peers rank Troy in the bottom third "
             "for IM-window recovery",
             font_size=14, color=GRAY_DARK)

    # Box 2 — LIGHT_RED
    b2x = b1x + box_w + gap
    add_rect(slide, b2x, box_y, box_w, box_h, LIGHT_RED)
    accent_bar(slide, b2x, box_y, box_w, ACCENT_RED)
    add_text(slide, b2x + Inches(0.15), box_y + Inches(0.15), box_w - Inches(0.3), Inches(0.6),
             "59 positions", font_size=30, bold=True, color=ACCENT_RED)
    add_text(slide, b2x + Inches(0.15), box_y + Inches(0.85), box_w - Inches(0.3), Inches(1.8),
             "Lost in national rankings — from #76 (top quarter) "
             "to #135 (mid-pack)",
             font_size=14, color=GRAY_DARK)

    # Box 3 — LIGHT_GREEN
    b3x = b2x + box_w + gap
    add_rect(slide, b3x, box_y, box_w + Inches(0.5), box_h, LIGHT_GREEN)
    accent_bar(slide, b3x, box_y, box_w + Inches(0.5), ACCENT_GREEN)
    add_text(slide, b3x + Inches(0.15), box_y + Inches(0.15), box_w, Inches(0.6),
             "9×", font_size=30, bold=True, color=ACCENT_GREEN)
    add_text(slide, b3x + Inches(0.15), box_y + Inches(0.85), box_w, Inches(1.8),
             "Birmingham recovered 9× more than Troy in the same "
             "window, same county, with tracking intact",
             font_size=14, color=GRAY_DARK)


# =============================================================================
# Slide 3 — What Troy Did and When
# =============================================================================
def build_slide_03():
    slide = new_slide(
        "What Troy Did and When",
        "Fall 2023: Illustrative Mathematics adopted, G6+G7 math detracked",
        slide_num=3)

    # Horizontal timeline
    line_y = Inches(3.5)
    line_left = Inches(1.2)
    line_right = Inches(12.0)
    # Horizontal line
    add_rect(slide, line_left, line_y, line_right - line_left, Inches(0.04), TROY_BLUE)

    events = [
        ("Fall 2023", "Adopted Illustrative Mathematics;\ndetracked G6 and G7 math", ACCENT_RED),
        ("Spring 2024", "First IM cohort tested\n(G6 M-STEP)", ACCENT_ORANGE),
        ("Fall 2024", "Second IM year; first G7 IM\ncohort enters 7th grade", ACCENT_ORANGE),
        ("Spring 2025", "Second IM G6 +\nfirst IM G7 tested", TROY_BLUE),
    ]

    spacing = (line_right - line_left) / (len(events) - 1)
    dot_size = Inches(0.22)

    for i, (period, desc, dot_color) in enumerate(events):
        cx = line_left + int(spacing) * i
        # Dot on line
        add_rect(slide, cx - dot_size // 2, line_y - dot_size // 2,
                 dot_size, dot_size, dot_color)
        # Period label above
        add_text(slide, cx - Inches(1.2), line_y - Inches(1.0), Inches(2.4), Inches(0.5),
                 period, font_size=16, bold=True, color=TROY_BLUE, alignment='center')
        # Description below
        add_text(slide, cx - Inches(1.3), line_y + Inches(0.3), Inches(2.6), Inches(1.2),
                 desc, font_size=13, color=GRAY_DARK, alignment='center')

    # Bottom note
    note_y = Inches(5.8)
    add_rect(slide, Inches(0.5), note_y, Inches(12.3), Inches(0.7), GRAY_LIGHT)
    add_text(slide, Inches(0.7), note_y + Inches(0.1), Inches(11.9), Inches(0.5),
             "Prior: Big Ideas Math 7, Math 7/8 Honors, small Algebra 1 cohort",
             font_size=13, color=GRAY_MID, valign='middle')


# =============================================================================
# Slide 4 — Pre-COVID Math Position
# =============================================================================
def build_slide_04():
    slide = new_slide(
        "Troy’s Pre-COVID Math Position — Strong, Not Invincible",
        "Pre-COVID G6+G7: +0.937 grade levels above national norm",
        slide_num=4)

    embed_chart(slide, "chart01_mi_peers_g6_trend.png")

    # Right panel — LIGHT_GREEN
    ix, iy, iw, ih = panel_box(slide, PANEL_L, PANEL_T, PANEL_W, PANEL_H, LIGHT_GREEN)
    add_text(slide, ix, iy, iw, Inches(0.7),
             "+0.937", font_size=32, bold=True, color=ACCENT_GREEN)
    add_multiline(slide, ix, iy + Inches(0.75), iw, ih - Inches(0.8), [
        "grade levels above national norm",
        "",
        "2nd highest among MI affluent peers",
        "",
        "Top quarter nationally (#76 / 296)",
    ], font_size=15, color=GRAY_DARK, space_after=6)


# =============================================================================
# Slide 5 — Troy Fell 59 Places
# =============================================================================
def build_slide_05():
    slide = new_slide(
        "Troy Fell 59 Places in the National Math Rankings",
        "From top-quarter to mid-pack among 296 level-matched peers",
        slide_num=5)

    embed_chart(slide, "chart11_rank_shift_scatter.png")

    ix, iy, iw, ih = panel_box(slide, PANEL_L, PANEL_T, PANEL_W, PANEL_H, LIGHT_RED)
    add_text(slide, ix, iy, iw, Inches(0.8),
             "#76 → #135", font_size=30, bold=True, color=ACCENT_RED)
    add_multiline(slide, ix, iy + Inches(0.85), iw, ih - Inches(0.9), [
        "Lost 59 positions",
        "",
        "Pre-COVID: top quarter",
        "Post-COVID: mid-pack",
        "",
        "Peers at the same starting level pulled ahead",
    ], font_size=15, color=GRAY_DARK, space_after=6)


# =============================================================================
# Slide 6 — MI Peer Leaderboard
# =============================================================================
def build_slide_06():
    slide = new_slide(
        "MI Peer Leaderboard: Troy Fell From 2nd to 4th",
        "4 of 5 subgroups lost positions under IM + detracking",
        slide_num=6)

    embed_chart(slide, "chart12_mi_peer_bump.png")

    # Right panel — LIGHT_RED with subgroup table
    ix, iy, iw, ih = panel_box(slide, PANEL_L, PANEL_T, PANEL_W, PANEL_H, LIGHT_RED)

    add_text(slide, ix, iy, iw, Inches(0.4),
             "Troy’s Rank Shift by Subgroup", font_size=14, bold=True, color=DARK_RED_HDR)

    subgroups = [
        ("All Students", "2nd→4th", "−2"),
        ("Asian", "1st→2nd", "−1"),
        ("White", "6th→5th", "+1"),
        ("Econ Disadv", "3rd→4th", "−1"),
        ("Not Econ Disadv", "2nd→3rd", "−1"),
    ]

    # Table header row
    row_h = Inches(0.42)
    tbl_y = iy + Inches(0.55)
    add_rect(slide, ix, tbl_y, iw, row_h, TROY_BLUE)
    col_w1, col_w2, col_w3 = Inches(1.85), Inches(1.25), Inches(0.7)
    add_text(slide, ix + Inches(0.08), tbl_y + Inches(0.05), col_w1, row_h,
             "Subgroup", font_size=11, bold=True, color=WHITE)
    add_text(slide, ix + col_w1, tbl_y + Inches(0.05), col_w2, row_h,
             "Shift", font_size=11, bold=True, color=WHITE)
    add_text(slide, ix + col_w1 + col_w2, tbl_y + Inches(0.05), col_w3, row_h,
             "Δ", font_size=11, bold=True, color=WHITE, alignment='center')

    for i, (sg, shift, delta) in enumerate(subgroups):
        ry = tbl_y + row_h + Inches(i * 0.42)
        bg = WHITE if i % 2 == 0 else GRAY_LIGHT
        add_rect(slide, ix, ry, iw, Inches(0.42), bg)
        clr = ACCENT_GREEN if delta.startswith("+") else ACCENT_RED
        add_text(slide, ix + Inches(0.08), ry + Inches(0.05), col_w1, Inches(0.35),
                 sg, font_size=12, bold=True, color=GRAY_DARK)
        add_text(slide, ix + col_w1, ry + Inches(0.05), col_w2, Inches(0.35),
                 shift, font_size=12, color=GRAY_DARK)
        add_text(slide, ix + col_w1 + col_w2, ry + Inches(0.05), col_w3, Inches(0.35),
                 delta, font_size=13, bold=True, color=clr, alignment='center')

    # Bottom note
    note_y = tbl_y + row_h + Inches(len(subgroups) * 0.42) + Inches(0.2)
    add_text(slide, ix, note_y, iw, Inches(0.7),
             "Was #1 Asian math in MI.\nNorthville now leads.",
             font_size=13, bold=True, color=ACCENT_RED)


# =============================================================================
# Slide 7 — Post-COVID Decline
# =============================================================================
def build_slide_07():
    slide = new_slide(
        "The Post-COVID Decline Hit Troy Hard in Math",
        "Troy declined −0.187 grade levels — rank 104 of 296 (35th percentile)",
        slide_num=7)

    embed_chart(slide, "chart09_covid_delta_ranking.png")

    ix, iy, iw, ih = panel_box(slide, PANEL_L, PANEL_T, PANEL_W, PANEL_H, LIGHT_RED)
    add_text(slide, ix, iy, iw, Inches(0.7),
             "−0.187", font_size=32, bold=True, color=ACCENT_RED)
    add_multiline(slide, ix, iy + Inches(0.75), iw, ih - Inches(0.8), [
        "grade-level decline",
        "",
        "Rank 104 of 296",
        "(35th percentile)",
        "",
        "3rd worst among 8 MI peers",
    ], font_size=15, color=GRAY_DARK, space_after=6)


# =============================================================================
# Slide 8 — The Critical Test: IM Window
# =============================================================================
def build_slide_08():
    slide = new_slide(
        "The Critical Test: Recovery During the IM Window (2022–23 → 2024–25)",
        "If IM were working, Troy should show stronger-than-peer recovery",
        slide_num=8)

    embed_chart(slide, "chart04_level_matched_histogram.png")

    ix, iy, iw, ih = panel_box(slide, PANEL_L, PANEL_T, PANEL_W, PANEL_H, LIGHT_RED)
    add_text(slide, ix, iy, iw, Inches(0.5),
             "Rank 197 / 295", font_size=28, bold=True, color=ACCENT_RED)
    add_text(slide, ix, iy + Inches(0.55), iw, Inches(0.4),
             "bottom third", font_size=18, bold=True, color=ACCENT_RED)
    add_multiline(slide, ix, iy + Inches(1.2), iw, ih - Inches(1.3), [
        "Troy: +0.024",
        "Peer median: +0.058",
        "",
        "Most peers recovered 2–3× more",
    ], font_size=15, color=GRAY_DARK, space_after=6)


# =============================================================================
# Slide 9 — MI Peers IM Window Ranking
# =============================================================================
def build_slide_09():
    slide = new_slide(
        "Among MI Peers, Troy Ranks 6th of 8 in Recovery",
        "Birmingham (tracked) recovered 9× more than Troy (IM + detracked)",
        slide_num=9)

    # Slightly narrower chart to fit the wider table panel
    embed_chart(slide, "chart03_im_window_ranking.png",
                left=CHART_L, top=CHART_T, width=Inches(7.0), height=CHART_H)

    # Wider panel for the ranking table
    tbl_l = Inches(7.5)
    tbl_w = Inches(5.5)
    ix, iy, iw, ih = panel_box(slide, tbl_l, PANEL_T, tbl_w, PANEL_H, GRAY_LIGHT)

    add_text(slide, ix, iy, iw, Inches(0.4),
             "IM-Window Recovery Ranking", font_size=14, bold=True, color=TROY_BLUE)

    peers = [
        ("Birmingham", "+0.211", "(tracked)", ACCENT_GREEN),
        ("Northville", "+0.102", "(tracked)", ACCENT_GREEN),
        ("Bloomfield Hills", "+0.082", "(tracked)", ACCENT_GREEN),
        ("Walled Lake", "+0.056", "", GRAY_DARK),
        ("Novi", "+0.029", "", GRAY_DARK),
        ("Troy", "+0.024", "(IM+detracked)", ACCENT_RED),
        ("Rochester", "+0.000", "", GRAY_DARK),
        ("West Bloomfield", "+0.006", "", GRAY_DARK),
    ]

    row_h_val = Inches(0.56)
    for i, (name, delta, note, clr) in enumerate(peers):
        ry = iy + Inches(0.5) + int(row_h_val) * i
        is_troy = (name == "Troy")
        if is_troy:
            add_rect(slide, ix, ry - Inches(0.02), iw, row_h_val, LIGHT_RED)
        add_text(slide, ix + Inches(0.1), ry + Inches(0.06), Inches(2.0), Inches(0.4),
                 name, font_size=13, bold=is_troy, color=clr)
        add_text(slide, ix + Inches(2.2), ry + Inches(0.06), Inches(1.1), Inches(0.4),
                 delta, font_size=13, bold=True, color=clr)
        if note:
            add_text(slide, ix + Inches(3.4), ry + Inches(0.06), Inches(1.8), Inches(0.4),
                     note, font_size=11, color=GRAY_MID)


# =============================================================================
# Slide 10 — G7 Tracked vs IM
# =============================================================================
def build_slide_10():
    slide = new_slide(
        "G7 Math: The Tracked-vs-IM Comparison Year",
        "2023-24: G7 tracked. 2024-25: First G7 IM cohort — recovery stalled",
        slide_num=10)

    embed_chart(slide, "chart02_mi_peers_g7_trend.png")

    # Two stacked panels
    half_h = Inches(2.55)
    gap = Inches(0.15)

    # Top panel — LIGHT_GREEN
    ix, iy, iw, ih = panel_box(slide, PANEL_L, PANEL_T, PANEL_W, half_h, LIGHT_GREEN)
    add_text(slide, ix, iy, iw, Inches(0.35),
             "2023-24: G7 still tracked", font_size=14, bold=True, color=ACCENT_GREEN)
    add_text(slide, ix, iy + Inches(0.45), iw, ih - Inches(0.5),
             "Troy matched peers when G7 still used the tracked curriculum "
             "with Big Ideas Math and honors pathways.",
             font_size=14, color=GRAY_DARK)

    # Bottom panel — LIGHT_RED
    p2_top = PANEL_T + half_h + gap
    ix2, iy2, iw2, ih2 = panel_box(slide, PANEL_L, p2_top, PANEL_W, half_h + Inches(0.25), LIGHT_RED)
    add_text(slide, ix2, iy2, iw2, Inches(0.35),
             "2024-25: First G7 IM cohort", font_size=14, bold=True, color=ACCENT_RED)
    add_text(slide, ix2, iy2 + Inches(0.45), iw2, ih2 - Inches(0.5),
             "Recovery stalled. The first cohort to experience IM in both "
             "G6 and G7 shows no additional gains over the prior tracked year.",
             font_size=14, color=GRAY_DARK)


# =============================================================================
# Slide 11 — Asian Subgroup Crisis
# =============================================================================
def build_slide_11():
    slide = new_slide(
        "The Asian Subgroup Crisis — 36% of Troy’s Enrollment",
        "Troy’s historically strongest demographic shows the sharpest decline",
        slide_num=11)

    embed_chart(slide, "chart05_asian_g6_trend.png")

    ix, iy, iw, ih = panel_box(slide, PANEL_L, PANEL_T, PANEL_W, PANEL_H, LIGHT_RED)

    stats = [
        ("G6 Asian Δ COVID:", "−0.409", ACCENT_RED),
        ("G7 Asian Δ COVID:", "−0.278", ACCENT_RED),
        ("G7 Asian IM-window:", "−0.094", ACCENT_RED),
    ]

    for i, (label, value, clr) in enumerate(stats):
        sy = iy + Inches(i * 1.6)
        add_text(slide, ix, sy, iw, Inches(0.35),
                 label, font_size=14, bold=True, color=GRAY_DARK)
        add_text(slide, ix, sy + Inches(0.4), iw, Inches(0.6),
                 value, font_size=30, bold=True, color=clr)
        if i == 2:
            add_text(slide, ix, sy + Inches(1.05), iw, Inches(0.35),
                     "(still declining)", font_size=14, bold=True, color=ACCENT_RED)


# =============================================================================
# Slide 12 — G7 Asian Still Declining
# =============================================================================
def build_slide_12():
    slide = new_slide(
        "G7 Asian Math Is Still Declining Under IM",
        "High-Asian peers with tracking show robust recovery",
        slide_num=12)

    embed_chart(slide, "chart06_asian_g7_trend.png")

    # Comparison cards
    cards = [
        ("Bellevue (WA)", "+0.162", LIGHT_GREEN, ACCENT_GREEN),
        ("Issaquah (WA)", "+0.218", LIGHT_GREEN, ACCENT_GREEN),
        ("Troy (MI)", "−0.094", LIGHT_RED, ACCENT_RED),
    ]

    card_h = Inches(1.55)
    gap = Inches(0.15)
    for i, (dist, delta, bg, clr) in enumerate(cards):
        cy = PANEL_T + Inches(i * (1.55 + 0.15))
        ix, iy, iw, ih = panel_box(slide, PANEL_L, cy, PANEL_W, card_h, bg)
        add_text(slide, ix, iy, iw, Inches(0.35),
                 dist, font_size=14, bold=True, color=GRAY_DARK)
        add_text(slide, ix, iy + Inches(0.45), iw, Inches(0.7),
                 delta, font_size=30, bold=True, color=clr)


# =============================================================================
# Slide 13 — Every Subgroup Declined
# =============================================================================
def build_slide_13():
    slide = new_slide(
        "Every Subgroup Declined — But Recovery Is Uneven",
        "Asian and White subgroups show largest absolute declines",
        slide_num=13)

    embed_chart(slide, "chart07_troy_subgroup_waterfall.png")

    ix, iy, iw, ih = panel_box(slide, PANEL_L, PANEL_T, PANEL_W, PANEL_H, GRAY_LIGHT)
    add_text(slide, ix, iy, iw, Inches(0.4),
             "Key Findings", font_size=16, bold=True, color=TROY_BLUE)
    add_multiline(slide, ix, iy + Inches(0.5), iw, ih - Inches(0.6), [
        "Asian and White subgroups show the largest absolute declines from pre-COVID levels.",
        "",
        "Economically disadvantaged students were already near the national average — "
        "the gap between ECD and non-ECD has not narrowed.",
        "",
        "No subgroup has returned to pre-COVID performance under IM + detracking.",
        "",
        "The “rising tide lifts all boats” promise of detracking has not materialised "
        "for any demographic group in Troy.",
    ], font_size=14, color=GRAY_DARK, space_after=4)


# =============================================================================
# Slide 14 — M-STEP Proficiency
# =============================================================================
def build_slide_14():
    slide = new_slide(
        "M-STEP Math Proficiency: The Numbers Parents See",
        "Michigan M-STEP % proficient or advanced, grades 6 and 7",
        slide_num=14)

    embed_chart(slide, "chart08_mstep_troy_proficiency.png")

    ix, iy, iw, ih = panel_box(slide, PANEL_L, PANEL_T, PANEL_W, PANEL_H, GRAY_LIGHT)

    add_text(slide, ix, iy, iw, Inches(0.4),
             "% Proficient or Advanced", font_size=15, bold=True, color=TROY_BLUE)

    # Table header
    tbl_y = iy + Inches(0.55)
    row_h = Inches(0.38)
    add_rect(slide, ix, tbl_y, iw, row_h, TROY_BLUE)
    add_text(slide, ix + Inches(0.08), tbl_y + Inches(0.04), Inches(0.6), row_h,
             "Grade", font_size=11, bold=True, color=WHITE)
    years = ["2021-22", "2022-23", "2023-24", "2024-25"]
    col_w = Inches(0.85)
    for j, yr in enumerate(years):
        add_text(slide, ix + Inches(0.7) + col_w * j, tbl_y + Inches(0.04), col_w, row_h,
                 yr, font_size=10, bold=True, color=WHITE, alignment='center')

    # G6 row
    g6_y = tbl_y + row_h
    add_rect(slide, ix, g6_y, iw, row_h, WHITE)
    add_text(slide, ix + Inches(0.08), g6_y + Inches(0.04), Inches(0.6), row_h,
             "G6", font_size=12, bold=True, color=GRAY_DARK)
    g6_vals = ["71.4%", "63.2%", "66.3%", "65.8%"]
    for j, v in enumerate(g6_vals):
        add_text(slide, ix + Inches(0.7) + col_w * j, g6_y + Inches(0.04), col_w, row_h,
                 v, font_size=12, color=GRAY_DARK, alignment='center')

    # G7 row
    g7_y = g6_y + row_h
    add_rect(slide, ix, g7_y, iw, row_h, GRAY_LIGHT)
    add_text(slide, ix + Inches(0.08), g7_y + Inches(0.04), Inches(0.6), row_h,
             "G7", font_size=12, bold=True, color=GRAY_DARK)
    g7_vals = ["67.5%", "61.3%", "66.0%", "TBD"]
    for j, v in enumerate(g7_vals):
        add_text(slide, ix + Inches(0.7) + col_w * j, g7_y + Inches(0.04), col_w, row_h,
                 v, font_size=12, color=GRAY_DARK, alignment='center')

    # Bottom callout
    callout_y = g7_y + row_h + Inches(0.3)
    add_rect(slide, ix, callout_y, iw, Inches(1.2), LIGHT_RED)
    accent_bar(slide, ix, callout_y, iw, ACCENT_RED)
    add_text(slide, ix + Inches(0.1), callout_y + Inches(0.15), iw - Inches(0.2), Inches(1.0),
             "IM has not restored Troy to pre-COVID proficiency",
             font_size=14, bold=True, color=ACCENT_RED)


# =============================================================================
# Slide 15 — San Francisco
# =============================================================================
def build_slide_15():
    slide = new_slide(
        "San Francisco: 12 Years of Detracking, Then Reversal",
        "SFUSD eliminated Algebra 1 in 8th grade in 2014 — reversed March 2026",
        slide_num=15)

    findings = [
        (LIGHT_RED, ACCENT_RED,
         "1", "Achievement gaps WIDENED: Hispanic-White gap expanded "
         "31 points (vs 5 statewide)"),
        (LIGHT_ORANGE, ACCENT_ORANGE,
         "2", "AP math enrollment DECLINED 15%, driven by "
         "Asian/Pacific Islander students"),
        (LIGHT_RED, ACCENT_RED,
         "3", "Board voted March 2026 to RESTORE Algebra 1 in "
         "8th grade at all schools"),
    ]

    box_y = Inches(1.15)
    box_h = Inches(1.4)
    gap = Inches(0.15)

    for i, (bg, bar_clr, num, text) in enumerate(findings):
        ry = box_y + Inches(i * (1.4 + 0.15))
        add_rect(slide, Inches(0.4), ry, Inches(12.5), box_h, bg)
        accent_bar(slide, Inches(0.4), ry, Inches(12.5), bar_clr)
        # Number circle
        cx = Inches(0.7)
        cy_c = ry + Inches(0.35)
        add_rect(slide, cx, cy_c, Inches(0.45), Inches(0.45), bar_clr)
        add_text(slide, cx, cy_c, Inches(0.45), Inches(0.45),
                 num, font_size=16, bold=True, color=WHITE, alignment='center', valign='middle')
        # Text
        add_text(slide, Inches(1.35), ry + Inches(0.2), Inches(11.3), box_h - Inches(0.3),
                 text, font_size=16, color=GRAY_DARK)

    # Bottom callout bar
    callout_y = box_y + Inches(3 * (1.4 + 0.15)) + Inches(0.15)
    add_rect(slide, Inches(0.4), callout_y, Inches(12.5), Inches(0.8), TROY_BLUE)
    add_text(slide, Inches(0.7), callout_y + Inches(0.1), Inches(12.0), Inches(0.6),
             "Troy is replicating a strategy with a well-documented failure pattern",
             font_size=16, bold=True, color=WHITE, valign='middle')


# =============================================================================
# Slide 16 — Cambridge
# =============================================================================
def build_slide_16():
    slide = new_slide(
        "Cambridge, MA: Detracked 2017–2019, Now Reversing",
        "A second major detracking experiment now being reversed",
        slide_num=16)

    col_w = Inches(6.1)
    col_h = Inches(5.5)
    gap = Inches(0.3)
    y = Inches(1.1)

    # Left panel — LIGHT_ORANGE
    x1 = Inches(0.4)
    ix1, iy1, iw1, ih1 = panel_box(slide, x1, y, col_w, col_h, LIGHT_ORANGE)
    add_rect(slide, ix1, iy1, iw1, Inches(0.45), ACCENT_ORANGE)
    add_text(slide, ix1 + Inches(0.1), iy1 + Inches(0.05), iw1 - Inches(0.2), Inches(0.4),
             "What They Did", font_size=16, bold=True, color=WHITE)
    add_multiline(slide, ix1, iy1 + Inches(0.65), iw1, ih1 - Inches(0.7), [
        "Eliminated accelerated math tracks across all "
        "middle schools between 2017 and 2019.",
        "",
        "Goal: address racial disparities in advanced "
        "course enrollment by placing all students in "
        "the same heterogeneous math classes.",
        "",
        "Mirrored SFUSD’s rationale: equal access "
        "would close gaps.",
    ], font_size=15, color=GRAY_DARK, space_after=4)

    # Right panel — LIGHT_RED
    x2 = x1 + col_w + gap
    ix2, iy2, iw2, ih2 = panel_box(slide, x2, y, col_w, col_h, LIGHT_RED)
    add_rect(slide, ix2, iy2, iw2, Inches(0.45), ACCENT_RED)
    add_text(slide, ix2 + Inches(0.1), iy2 + Inches(0.05), iw2 - Inches(0.2), Inches(0.4),
             "What Happened", font_size=16, bold=True, color=WHITE)
    add_multiline(slide, ix2, iy2 + Inches(0.65), iw2, ih2 - Inches(0.7), [
        "COVID obscured early results, but post-pandemic "
        "data showed persistent achievement gaps.",
        "",
        "District is now reversing course and restoring "
        "algebra pathways in 8th grade.",
        "",
        "Advanced students left for private schools "
        "or after-school tutoring — exacerbating the "
        "very inequities the policy aimed to fix.",
    ], font_size=15, color=GRAY_DARK, space_after=4)


# =============================================================================
# Slide 17 — The 50% Algebra 1 Claim
# =============================================================================
def build_slide_17():
    slide = new_slide(
        "The ‘50% Algebra 1’ Claim — Course Access ≠ Achievement",
        "Course enrollment is not an achievement metric",
        slide_num=17)

    findings = [
        ("1", "SFUSD made the identical claim for a decade before test scores "
         "revealed no improvement"),
        ("2", "SEDA shows Troy’s G6+G7 math readiness dropped from "
         "+0.937 to +0.750 — a quarter-grade-level decline"),
        ("3", "The comparison is not past Troy enrollment — it’s what "
         "tracked peers achieve now"),
    ]

    box_y = Inches(1.15)
    box_h = Inches(1.55)
    gap = Inches(0.15)

    for i, (num, text) in enumerate(findings):
        ry = box_y + Inches(i * (1.55 + 0.15))
        add_rect(slide, Inches(0.4), ry, Inches(12.5), box_h, LIGHT_ORANGE)
        accent_bar(slide, Inches(0.4), ry, Inches(12.5), ACCENT_ORANGE)
        # Number badge
        cx = Inches(0.7)
        cy_c = ry + Inches(0.4)
        add_rect(slide, cx, cy_c, Inches(0.45), Inches(0.45), ACCENT_ORANGE)
        add_text(slide, cx, cy_c, Inches(0.45), Inches(0.45),
                 num, font_size=16, bold=True, color=WHITE, alignment='center', valign='middle')
        add_text(slide, Inches(1.35), ry + Inches(0.25), Inches(11.3), box_h - Inches(0.4),
                 text, font_size=16, color=GRAY_DARK)


# =============================================================================
# Slide 18 — What High-Performing Districts Do
# =============================================================================
def build_slide_18():
    slide = new_slide(
        "What High-Performing Math Districts Do Differently",
        "Among 107 high-Asian (≥20%) peers, many recovered with tracking intact",
        slide_num=18)

    embed_chart(slide, "chart10_high_asian_peers_delta.png")

    ix, iy, iw, ih = panel_box(slide, PANEL_L, PANEL_T, PANEL_W, PANEL_H, LIGHT_GREEN)
    add_text(slide, ix, iy, iw, Inches(0.4),
             "Key Finding", font_size=16, bold=True, color=ACCENT_GREEN)
    add_multiline(slide, ix, iy + Inches(0.5), iw, ih - Inches(0.6), [
        "Among 107 high-Asian (≥20%) districts nationally, "
        "many California and New Jersey districts recovered "
        "to or surpassed pre-COVID levels.",
        "",
        "These districts maintained tracking and traditional "
        "acceleration pathways.",
        "",
        "Troy is an outlier among its demographic peers — "
        "one of the few to adopt detracking.",
        "",
        "The data does not support the claim that detracking "
        "is necessary for equitable outcomes.",
    ], font_size=14, color=GRAY_DARK, space_after=4)


# =============================================================================
# Slide 19 — IM Evidence Gaps
# =============================================================================
def build_slide_19():
    slide = new_slide(
        "Illustrative Mathematics: Evidence Gaps",
        "No published evidence of IM success in a district like Troy",
        slide_num=19)

    # Table header
    tbl_y = Inches(1.15)
    tbl_x = Inches(0.4)
    tbl_w = Inches(12.5)
    row_h = Inches(0.5)

    col_x = [tbl_x + Inches(0.15), tbl_x + Inches(3.8), tbl_x + Inches(7.0)]
    col_w_vals = [Inches(3.5), Inches(3.0), Inches(5.3)]

    add_rect(slide, tbl_x, tbl_y, tbl_w, row_h, TROY_BLUE)
    headers = ["District", "Result", "Comparability"]
    for i, h in enumerate(headers):
        add_text(slide, col_x[i], tbl_y + Inches(0.05), col_w_vals[i], row_h,
                 h, font_size=14, bold=True, color=WHITE)

    table_data = [
        ("NYC District 11 (Bronx)", "25.8% → 50.6%",
         "High-poverty, low-baseline; not comparable to Troy"),
        ("Fort Zumwalt (MO)", "Effect size 0.16 SD",
         "Modest improvement; comparable to general instructional improvement"),
        ("Philadelphia", "Process study only",
         "No student outcomes data published"),
    ]

    for i, (district, result, note) in enumerate(table_data):
        ry = tbl_y + row_h + Inches(i * 0.95)
        bg = WHITE if i % 2 == 0 else GRAY_LIGHT
        add_rect(slide, tbl_x, ry, tbl_w, Inches(0.9), bg)
        add_text(slide, col_x[0], ry + Inches(0.15), col_w_vals[0], Inches(0.6),
                 district, font_size=14, bold=True, color=GRAY_DARK)
        add_text(slide, col_x[1], ry + Inches(0.15), col_w_vals[1], Inches(0.6),
                 result, font_size=14, color=GRAY_DARK)
        add_text(slide, col_x[2], ry + Inches(0.15), col_w_vals[2], Inches(0.6),
                 note, font_size=13, color=GRAY_MID)

    # Bottom callout
    callout_y = tbl_y + row_h + Inches(3 * 0.95) + Inches(0.2)
    add_rect(slide, tbl_x, callout_y, tbl_w, Inches(1.0), LIGHT_ORANGE)
    accent_bar(slide, tbl_x, callout_y, tbl_w, ACCENT_ORANGE)
    add_text(slide, tbl_x + Inches(0.2), callout_y + Inches(0.15), tbl_w - Inches(0.4), Inches(0.7),
             "No published evidence of IM success in an affluent, high-performing, "
             "high-Asian district like Troy. The existing studies all involve "
             "high-poverty, low-baseline populations.",
             font_size=14, bold=True, color=ACCENT_ORANGE)


# =============================================================================
# Slide 20 — Conclusions
# =============================================================================
def build_slide_20():
    slide = new_slide(
        "Conclusions",
        "Five findings from three independent lines of evidence",
        slide_num=20)

    conclusions = [
        (ACCENT_RED, LIGHT_RED,
         "Troy has NOT outperformed — bottom third of 295 peers"),
        (ACCENT_RED, LIGHT_RED,
         "Asian subgroup most concerning — G7 Asian STILL declining"),
        (ACCENT_ORANGE, LIGHT_ORANGE,
         "Detracking evidence base against — SFUSD, Cambridge reversed"),
        (ACCENT_ORANGE, LIGHT_ORANGE,
         "Course enrollment ≠ achievement — 50% Algebra 1 misleads"),
        (ACCENT_RED, LIGHT_RED,
         "Opportunity cost real — Birmingham recovered 9× more"),
    ]

    y = Inches(1.15)
    row_h = Inches(1.05)
    gap = Inches(0.12)

    for i, (circle_clr, bg_clr, text) in enumerate(conclusions):
        ry = y + Inches(i * (1.05 + 0.12))
        # Background
        add_rect(slide, Inches(0.5), ry, Inches(12.3), row_h, bg_clr)
        # Colored circle (square as approximation)
        cx = Inches(0.7)
        cy = ry + Inches(0.22)
        add_rect(slide, cx, cy, Inches(0.5), Inches(0.5), circle_clr)
        add_text(slide, cx, cy, Inches(0.5), Inches(0.5),
                 str(i + 1), font_size=18, bold=True, color=WHITE,
                 alignment='center', valign='middle')
        # Left accent bar
        add_rect(slide, Inches(0.5), ry, Inches(0.06), row_h, circle_clr)
        # Text
        add_text(slide, Inches(1.4), ry + Inches(0.1), Inches(11.2), row_h - Inches(0.2),
                 text, font_size=16, bold=False, color=GRAY_DARK, valign='middle')


# =============================================================================
# Slide 21 — Methodology
# =============================================================================
def build_slide_21():
    slide = new_slide(
        "Methodology",
        "SEDA 2025.1 (Stanford CEPA) — NAEP-anchored cross-state comparison",
        slide_num=21)

    sections = [
        ("Primary Metric",
         "SEDA cs (cohort-standardized), NAEP-anchored, expressed in grade levels "
         "above/below the national average. Enables cross-state, cross-year comparison."),
        ("Time Windows",
         "Pre-COVID baseline: 2017–19 mean. Post-COVID: 2022–25 mean. "
         "IM window: 2022–23 → 2024–25 (the exact adoption period)."),
        ("Peer Groups",
         "296 level-matched peers: pre-COVID G6+G7 math ±0.25 of Troy’s +0.937. "
         "107 high-Asian peers: ≥20% Asian enrollment, pre-COVID ≥+0.50. "
         "8 MI affluent peers: same state test, comparable demographics."),
        ("M-STEP Data",
         "Michigan School Data (mischooldata.org). % proficient or advanced, "
         "grades 6 and 7, 2021–22 through 2024–25."),
    ]

    y = Inches(1.15)
    sec_h = Inches(1.3)
    gap = Inches(0.1)

    for i, (header, body) in enumerate(sections):
        sy = y + Inches(i * (1.3 + 0.1))
        # Header bar
        add_rect(slide, Inches(0.4), sy, Inches(12.5), Inches(0.4), TROY_BLUE)
        add_text(slide, Inches(0.6), sy + Inches(0.04), Inches(12.0), Inches(0.35),
                 header, font_size=13, bold=True, color=WHITE)
        # Body
        add_rect(slide, Inches(0.4), sy + Inches(0.4), Inches(12.5), sec_h - Inches(0.4), WHITE)
        add_text(slide, Inches(0.6), sy + Inches(0.45), Inches(12.0), sec_h - Inches(0.5),
                 body, font_size=13, color=GRAY_DARK)


# =============================================================================
# Slide 22 — References
# =============================================================================
def build_slide_22():
    slide = new_slide(
        "References",
        "Primary data sources and research",
        slide_num=22)

    references = [
        "SEDA 2025.1 — Stanford Center for Education Policy Analysis (CEPA). "
        "educationdata.urban.org / edopportunity.org",
        "Education Scorecard 2026 — Harvard-Stanford-Dartmouth collaboration. "
        "educationscorecard.org",
        "MI School Data / M-STEP — Michigan Department of Education. "
        "mischooldata.org",
        "Dee, T., Huffaker, E., & Novicoff, S. (2025). “Detracking and "
        "Advanced Course-Taking in San Francisco.” SFUSD study.",
        "Education Next — SFUSD detracking analysis and coverage of the "
        "March 2026 board reversal vote.",
        "Brookings Institution — Detracking evidence review. Comprehensive "
        "survey of heterogeneous grouping outcomes.",
    ]

    y = Inches(1.15)
    row_h = Inches(0.88)

    for i, ref in enumerate(references):
        ry = y + Inches(i * 0.92)
        bg = WHITE if i % 2 == 0 else GRAY_LIGHT
        add_rect(slide, Inches(0.4), ry, Inches(12.5), row_h, bg)
        # Number
        add_text(slide, Inches(0.55), ry + Inches(0.1), Inches(0.5), row_h - Inches(0.15),
                 f"{i + 1}.", font_size=15, bold=True, color=TROY_BLUE)
        # Reference text
        add_text(slide, Inches(1.0), ry + Inches(0.1), Inches(11.7), row_h - Inches(0.15),
                 ref, font_size=13, color=GRAY_DARK)


# =============================================================================
# Build all 22 slides
# =============================================================================
def main():
    build_slide_01()
    build_slide_02()
    build_slide_03()
    build_slide_04()
    build_slide_05()
    build_slide_06()
    build_slide_07()
    build_slide_08()
    build_slide_09()
    build_slide_10()
    build_slide_11()
    build_slide_12()
    build_slide_13()
    build_slide_14()
    build_slide_15()
    build_slide_16()
    build_slide_17()
    build_slide_18()
    build_slide_19()
    build_slide_20()
    build_slide_21()
    build_slide_22()

    prs.save(str(OUTPUT))
    print(f"Deck saved to {OUTPUT}")
    print(f"  {len(prs.slides)} slides, {SLIDE_W} x {SLIDE_H}")
    size_kb = OUTPUT.stat().st_size / 1024
    print(f"  File size: {size_kb:.0f} KB")


if __name__ == "__main__":
    main()
